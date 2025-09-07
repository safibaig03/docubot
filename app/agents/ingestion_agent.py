import os
import pandas as pd
import docx
from pptx import Presentation
from pypdf import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from typing import List, Dict, Any

# This agent is responsible for all document processing tasks.
# It takes raw uploaded files, parses them into text, and splits them into chunks.
class IngestionAgent:
    def __init__(self):
        # Initialize the text splitter. This will be used to break down large texts.
        # It recursively tries to split on paragraphs, then sentences, then words.
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100
        )

    def process_files(self, files: List[Any]) -> List[Dict[str, Any]]:
        # This is the main public method for the agent. It orchestrates the entire ingestion pipeline.
        all_docs = []
        # First, parse all uploaded files to extract text and initial metadata.
        for file in files:
            all_docs.extend(self._parse_document(file))
            
        # Second, chunk the parsed documents into smaller pieces.
        final_chunks = []
        doc_id_counter = 0
        for doc in all_docs:
            split_chunks = self.text_splitter.split_text(doc["content"])
            for chunk in split_chunks:
                # Ensure each chunk retains the metadata of its parent document.
                chunk_metadata = doc["metadata"].copy()
                # Create a unique ID for each individual chunk for the vector database.
                chunk_id = f"{doc['metadata']['source']}-chunk-{doc_id_counter}"
                
                final_chunks.append({
                    "id": chunk_id,
                    "content": chunk,
                    "metadata": chunk_metadata
                })
                doc_id_counter += 1
        return final_chunks

    def _parse_document(self, file: Any) -> List[Dict[str, Any]]:
        # This private method acts as a router, calling the correct parser based on file extension.
        file_extension = os.path.splitext(file.filename)[1].lower()
        
        if file_extension == ".pdf":
            return self._parse_pdf(file)
        elif file_extension == ".docx":
            return self._parse_docx(file)
        elif file_extension == ".pptx":
            return self._parse_pptx(file)
        elif file_extension == ".csv":
            return self._parse_csv(file)
        elif file_extension in [".txt", ".md"]:
            return self._parse_txt(file)
        else:
            return []

    # --- Specialized Parsing Methods ---
    # Each of the following methods handles a specific file type.
    # They all return a list of dictionaries, each containing 'content' and 'metadata'.

    def _parse_pdf(self, file: Any) -> List[Dict[str, Any]]:
        reader = PdfReader(file.file)
        chunks = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text: chunks.append({"content": text, "metadata": {"source": file.filename, "page": i + 1}})
        return chunks

    def _parse_docx(self, file: Any) -> List[Dict[str, Any]]:
        document = docx.Document(file.file)
        chunks = []
        for i, para in enumerate(document.paragraphs):
            if para.text.strip(): chunks.append({"content": para.text, "metadata": {"source": file.filename, "paragraph": i + 1}})
        return chunks

    def _parse_pptx(self, file: Any) -> List[Dict[str, Any]]:
        pres = Presentation(file.file)
        chunks = []
        for slide_num, slide in enumerate(pres.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"): slide_text += shape.text + "\n"
            if slide_text.strip(): chunks.append({"content": slide_text, "metadata": {"source": file.filename, "slide": slide_num + 1}})
        return chunks

    def _parse_csv(self, file: Any) -> List[Dict[str, Any]]:
        df = pd.read_csv(file.file)
        chunks = []
        # Convert each row of the CSV into a single text string.
        for i, row in df.iterrows():
            row_text = ", ".join(f"{col}: {val}" for col, val in row.items())
            chunks.append({"content": row_text, "metadata": {"source": file.filename, "row": i + 1}})
        return chunks

    def _parse_txt(self, file: Any) -> List[Dict[str, Any]]:
        content = file.file.read().decode('utf-8')
        return [{"content": content, "metadata": {"source": file.filename}}]