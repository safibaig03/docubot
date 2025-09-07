import os
import pandas as pd
import docx
from pptx import Presentation
from pypdf import PdfReader
import chromadb
from chromadb.utils import embedding_functions
import google.generativeai as genai
from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter

load_dotenv()
genai.configure(api_key=os.environ["GOOGLE_API_KEY"])

def parse_document(file):
    file_extension = os.path.splitext(file.name)[1].lower()
    if file_extension == ".pdf": return _parse_pdf(file)
    elif file_extension == ".docx": return _parse_docx(file)
    elif file_extension == ".pptx": return _parse_pptx(file)
    elif file_extension == ".csv": return _parse_csv(file)
    elif file_extension in [".txt", ".md"]: return _parse_txt(file)
    else: return []

def _parse_pdf(file):
    reader = PdfReader(file)
    chunks = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text: chunks.append({"content": text, "metadata": {"source": file.name, "page": i + 1}})
    return chunks

def _parse_docx(file):
    document = docx.Document(file)
    chunks = []
    for i, para in enumerate(document.paragraphs):
        if para.text.strip(): chunks.append({"content": para.text, "metadata": {"source": file.name, "paragraph": i + 1}})
    return chunks

def _parse_pptx(file):
    pres = Presentation(file)
    chunks = []
    for slide_num, slide in enumerate(pres.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"): slide_text += shape.text + "\n"
        if slide_text.strip(): chunks.append({"content": slide_text, "metadata": {"source": file.name, "slide": slide_num + 1}})
    return chunks

def _parse_csv(file):
    df = pd.read_csv(file)
    chunks = []
    for i, row in df.iterrows():
        row_text = ", ".join(f"{col}: {val}" for col, val in row.items())
        chunks.append({"content": row_text, "metadata": {"source": file.name, "row": i + 1}})
    return chunks

def _parse_txt(file):
    content = file.read().decode('utf-8')
    return [{"content": content, "metadata": {"source": file.name}}]

class RAGCore:
    def __init__(self):
        self.llm = genai.GenerativeModel('gemini-1.5-flash')
        self.embedding_function = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
        client = chromadb.Client()
        self.collection = client.get_or_create_collection(name="document_collection", embedding_function=self.embedding_function)
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)

    def add_documents(self, uploaded_files):
        all_docs = []
        for file in uploaded_files:
            all_docs.extend(parse_document(file))
        final_chunks = []
        doc_id_counter = 0
        for doc in all_docs:
            split_chunks = self.text_splitter.split_text(doc["content"])
            for chunk in split_chunks:
                chunk_metadata = doc["metadata"].copy()
                chunk_id = f"{doc['metadata']['source']}-chunk-{doc_id_counter}"
                final_chunks.append({"id": chunk_id, "content": chunk, "metadata": chunk_metadata})
                doc_id_counter += 1
        if not final_chunks: return
        self.collection.add(ids=[c['id'] for c in final_chunks], documents=[c['content'] for c in final_chunks], metadatas=[c['metadata'] for c in final_chunks])

    def get_answer(self, query):
        results = self.collection.query(
            query_texts=[query],
            n_results=3,
            include=["metadatas", "documents"] # Explicitly request metadatas and documents
        )
                
        retrieved_docs_content = results.get('documents', [[]])[0]
        retrieved_docs_metadata = results.get('metadatas', [[]])[0]
        
        sources = [
            {"content": content, "metadata": metadata}
            for content, metadata in zip(retrieved_docs_content, retrieved_docs_metadata)
        ]
        
        context_for_llm = ""
        for i, source in enumerate(sources):
            context_for_llm += f"Source {i+1} (from {source['metadata'].get('source', 'Unknown')}):\n"
            context_for_llm += source['content']
            context_for_llm += "\n\n---\n\n"

        prompt = f"""
        Answer the following question based only on the provided context from the sources.
        For each piece of information you use, you MUST cite the source filename from which it came.
        
        Context:
        {context_for_llm}
        
        Question:
        {query}
        """
        
        response = self.llm.generate_content(prompt)
        
        return {
            "answer": response.text,
            "sources": sources
        }